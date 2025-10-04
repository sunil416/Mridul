
from pptx import Presentation
from utils.Document import Document
class LoadPPTX:
    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self) -> list:
        prs = Presentation(self.file_path)
        documents = []

        for i, slide in enumerate(prs.slides):
            text = ""
            images = []

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
                if shape.shape_type == 13:  # Picture
                    image = shape.image
                    images.append(image.blob)

            documents.append(Document(
                page_content=text.strip(),
                metadata={"source": self.file_path, "slide": i + 1},
                resources=images  # ✅ fixed field name
            ))
        self._documents = documents
        return documents
    
